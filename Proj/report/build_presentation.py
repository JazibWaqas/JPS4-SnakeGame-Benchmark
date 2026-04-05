"""
Builds the 5-slide interim presentation.
Run: python3 build_presentation.py
Output: ADA_Interim_Presentation.pptx next to this script.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")

# Color palette (dark theme)
BG = RGBColor(0x0F, 0x17, 0x2A)        # deep navy
PANEL = RGBColor(0x1A, 0x25, 0x3D)     # lifted navy
CARD = RGBColor(0x22, 0x30, 0x4D)      # card
ACCENT = RGBColor(0x38, 0xBD, 0xF8)    # cyan
WHITE = RGBColor(0xEE, 0xF2, 0xF7)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
DIM = RGBColor(0x64, 0x74, 0x8B)
DIJ = RGBColor(0xE0, 0x5C, 0x5C)       # Dijkstra red
AST = RGBColor(0x4D, 0x94, 0xE0)       # A* blue
JPS = RGBColor(0x2C, 0xA0, 0x2C)       # JPS4 green
WARN = RGBColor(0xF5, 0xA6, 0x24)

# Slide dimensions (16:9)
SW, SH = 13.333, 7.5  # inches

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
blank = prs.slide_layouts[6]


def add_bg(slide, color=BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def add_rect(slide, x, y, w, h, fill=CARD, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.08
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size)
    run.font.bold = bold; run.font.color.rgb = color
    return tb


def add_accent_bar(slide, x, y, w=0.08, h=0.6, color=ACCENT):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def add_dot(slide, x, y, d=0.28, color=ACCENT):
    r = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def add_footer(slide, page, total=4):
    add_text(slide, 0.5, 7.1, 6, 0.3,
             "ADA Progress Report  |  Habib University  |  Spring 2026",
             size=9, color=DIM)
    add_text(slide, SW - 1.2, 7.1, 0.7, 0.3, f"{page} / {total}",
             size=9, color=DIM, align=PP_ALIGN.RIGHT)


# =============================================================================
# SLIDE 1: Title
# =============================================================================
s = prs.slides.add_slide(blank)
add_bg(s)

# Large decorative accent rectangle on the left
deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), prs.slide_height)
deco.fill.solid(); deco.fill.fore_color.rgb = ACCENT
deco.line.fill.background()
deco.shadow.inherit = False

# Hero image on the right (JPS concept)
jps_hero = os.path.join(FIG, "jps_concept.png")
if os.path.exists(jps_hero):
    pic = s.shapes.add_picture(jps_hero, Inches(9.0), Inches(2.1), height=Inches(3.4))

add_text(s, 0.9, 1.4, 11, 0.4, "HABIB UNIVERSITY  |  ALGORITHM DESIGN & ANALYSIS",
         size=11, bold=True, color=ACCENT)

add_text(s, 0.9, 1.9, 8.0, 1.4,
         "Comparing Shortest-Path",
         size=40, bold=True, color=WHITE)
add_text(s, 0.9, 2.7, 8.0, 1.0,
         "Algorithms on Grid Maps",
         size=40, bold=True, color=WHITE)
add_text(s, 0.9, 3.55, 8.0, 0.5,
         "Dijkstra vs A* vs JPS4 on 4-connected grids",
         size=16, color=MUTED)

# Algorithm chips
chips_y = 4.4
chip_w = 1.55
add_rect(s, 0.9, chips_y, chip_w, 0.55, fill=DIJ)
add_text(s, 0.9, chips_y, chip_w, 0.55, "Dijkstra", size=16, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 2.55, chips_y, chip_w, 0.55, fill=AST)
add_text(s, 2.55, chips_y, chip_w, 0.55, "A*", size=16, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 4.2, chips_y, chip_w, 0.55, fill=JPS)
add_text(s, 4.2, chips_y, chip_w, 0.55, "JPS4", size=16, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, 0.9, 5.4, 8, 0.45,
         "Interim Progress Report", size=18, bold=True, color=MUTED)

add_text(s, 0.9, 5.95, 8, 0.4,
         "Jazib Waqas    |    Salman Adnan    |    Hunain Abbas",
         size=14, color=WHITE)
add_text(s, 0.9, 6.35, 8, 0.35,
         "jw08048 . sa07885 . sh08466",
         size=11, color=DIM)

add_footer(s, 1)

# =============================================================================
# SLIDE 2: Paper Understanding / The Three Algorithms
# =============================================================================
s = prs.slides.add_slide(blank)
add_bg(s)
add_accent_bar(s, 0.5, 0.55, w=0.08, h=0.55, color=ACCENT)
add_text(s, 0.75, 0.5, 12, 0.65,
         "Three Algorithms, One Problem", size=28, bold=True, color=WHITE)
add_text(s, 0.75, 1.15, 12, 0.4,
         "Shortest path on a 4-connected grid. Same input, very different behavior.",
         size=13, color=MUTED)

# Three algorithm columns
col_y = 1.95
col_h = 4.5
col_w = 4.0
gap = 0.2
x0 = 0.5

# Dijkstra column
add_rect(s, x0, col_y, col_w, col_h, fill=CARD)
add_rect(s, x0, col_y, col_w, 0.7, fill=DIJ)
add_text(s, x0, col_y, col_w, 0.7, "DIJKSTRA", size=18, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x0 + 0.3, col_y + 0.9, col_w - 0.6, 0.4,
         "The baseline", size=13, bold=True, color=DIJ)
add_text(s, x0 + 0.3, col_y + 1.35, col_w - 0.6, 2.3,
         "Explores every reachable cell in order of distance from start. "
         "Always correct, but has no sense of where the goal is, so it "
         "essentially maps the entire reachable area first.",
         size=11, color=WHITE)
add_rect(s, x0 + 0.3, col_y + 3.5, col_w - 0.6, 0.8, fill=PANEL)
add_text(s, x0 + 0.3, col_y + 3.55, col_w - 0.6, 0.35,
         "Mechanism", size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, x0 + 0.3, col_y + 3.9, col_w - 0.6, 0.35,
         "Uniform wavefront", size=13, bold=True, color=DIJ, align=PP_ALIGN.CENTER)

# A* column
x1 = x0 + col_w + gap
add_rect(s, x1, col_y, col_w, col_h, fill=CARD)
add_rect(s, x1, col_y, col_w, 0.7, fill=AST)
add_text(s, x1, col_y, col_w, 0.7, "A*", size=18, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x1 + 0.3, col_y + 0.9, col_w - 0.6, 0.4,
         "Heuristic guided", size=13, bold=True, color=AST)
add_text(s, x1 + 0.3, col_y + 1.35, col_w - 0.6, 2.3,
         "Adds Manhattan distance to the goal as a lower-bound estimate. "
         "Picks the cell with the lowest total estimated cost at each step, "
         "so the search moves toward the goal instead of spreading evenly.",
         size=11, color=WHITE)
add_rect(s, x1 + 0.3, col_y + 3.5, col_w - 0.6, 0.8, fill=PANEL)
add_text(s, x1 + 0.3, col_y + 3.55, col_w - 0.6, 0.35,
         "Mechanism", size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, x1 + 0.3, col_y + 3.9, col_w - 0.6, 0.35,
         "Goal-biased search", size=13, bold=True, color=AST, align=PP_ALIGN.CENTER)

# JPS4 column
x2 = x1 + col_w + gap
add_rect(s, x2, col_y, col_w, col_h, fill=CARD)
add_rect(s, x2, col_y, col_w, 0.7, fill=JPS)
add_text(s, x2, col_y, col_w, 0.7, "JPS4", size=18, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x2 + 0.3, col_y + 0.9, col_w - 0.6, 0.4,
         "Baum (2025)", size=13, bold=True, color=JPS)
add_text(s, x2 + 0.3, col_y + 1.35, col_w - 0.6, 2.3,
         "Skips over long open corridors entirely. Stops only at jump points, "
         "places where an obstacle forces a real routing choice. Only those "
         "cells enter the open list, so the list stays small.",
         size=11, color=WHITE)
add_rect(s, x2 + 0.3, col_y + 3.5, col_w - 0.6, 0.8, fill=PANEL)
add_text(s, x2 + 0.3, col_y + 3.55, col_w - 0.6, 0.35,
         "Mechanism", size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, x2 + 0.3, col_y + 3.9, col_w - 0.6, 0.35,
         "Symmetry pruning", size=13, bold=True, color=JPS, align=PP_ALIGN.CENTER)

# Bottom insight strip
add_text(s, 0.75, 6.65, 12, 0.4,
         "All three guarantee the same shortest path. The difference is how many cells each one has to look at.",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)

add_footer(s, 2)

# =============================================================================
# SLIDE 3: Methodology
# =============================================================================
s = prs.slides.add_slide(blank)
add_bg(s)
add_accent_bar(s, 0.5, 0.55, w=0.08, h=0.55, color=ACCENT)
add_text(s, 0.75, 0.5, 12, 0.65,
         "Methodology", size=28, bold=True, color=WHITE)
add_text(s, 0.75, 1.15, 12, 0.4,
         "Same benchmark, same grids, same search loop. Only the algorithm swaps.",
         size=13, color=MUTED)

# Left: Setup parameters card
add_rect(s, 0.5, 1.9, 5.8, 5.0, fill=CARD)
add_text(s, 0.75, 2.05, 5.3, 0.45, "Benchmark Setup",
         size=16, bold=True, color=ACCENT)

setup_items = [
    ("Grid", "50 x 50 random obstacles"),
    ("Densities", "15 percent, 28 percent, 40 percent"),
    ("Trials per density", "24 accepted runs"),
    ("Seed", "42 (deterministic)"),
    ("Start and goal", "Placed far apart, verified reachable"),
    ("Acceptance rule", "Trial counted only if all three algorithms succeed"),
]
row_y = 2.7
for label, value in setup_items:
    add_text(s, 0.75, row_y, 2.0, 0.35, label.upper(),
             size=9, bold=True, color=MUTED)
    add_text(s, 0.75, row_y + 0.28, 5.3, 0.5, value,
             size=12, bold=True, color=WHITE)
    row_y += 0.7

# Right: Shared-loop architecture diagram
arch_x = 6.7
arch_w = 6.15
add_rect(s, arch_x, 1.9, arch_w, 5.0, fill=CARD)
add_text(s, arch_x + 0.25, 2.05, arch_w - 0.5, 0.45,
         "Shared Search Loop", size=16, bold=True, color=ACCENT)
add_text(s, arch_x + 0.25, 2.45, arch_w - 0.5, 0.4,
         "One A* engine. Heuristic and successor function plug in per algorithm.",
         size=10, color=MUTED)

# Central loop box
loop_x = arch_x + 1.85
loop_y = 3.1
loop_w = 2.4
loop_h = 0.85
add_rect(s, loop_x, loop_y, loop_w, loop_h, fill=ACCENT)
add_text(s, loop_x, loop_y, loop_w, loop_h, "astar_jps( )",
         size=16, bold=True, color=BG, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, loop_x, loop_y + loop_h + 0.05, loop_w, 0.3,
         "heap-based open list",
         size=9, color=MUTED, align=PP_ALIGN.CENTER)

# Three input boxes feeding into the loop
input_y = 5.05
bw = 1.85
bh = 1.4
bx = [arch_x + 0.3, arch_x + 2.15, arch_x + 4.0]
configs = [
    ("Dijkstra", DIJ, "h = 0", "4-neighbour successors"),
    ("A*", AST, "h = Manhattan", "4-neighbour successors"),
    ("JPS4", JPS, "h = Manhattan", "jump-point successors"),
]
for i, (name, color, h_text, s_text) in enumerate(configs):
    add_rect(s, bx[i], input_y, bw, bh, fill=PANEL, line=color)
    add_rect(s, bx[i], input_y, bw, 0.35, fill=color)
    add_text(s, bx[i], input_y, bw, 0.35, name, size=12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, bx[i] + 0.1, input_y + 0.42, bw - 0.2, 0.35,
             h_text, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, bx[i] + 0.1, input_y + 0.78, bw - 0.2, 0.55,
             s_text, size=9, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s, arch_x + 0.25, 6.55, arch_w - 0.5, 0.35,
         "Any measured difference comes from the algorithm itself.",
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_footer(s, 3)

# =============================================================================
# SLIDE 4: Preliminary Results
# =============================================================================
s = prs.slides.add_slide(blank)
add_bg(s)
add_accent_bar(s, 0.5, 0.55, w=0.08, h=0.55, color=ACCENT)
add_text(s, 0.75, 0.5, 12, 0.65,
         "Preliminary Results", size=28, bold=True, color=WHITE)
add_text(s, 0.75, 1.15, 12, 0.4,
         "Node expansions and wall-clock time vs. obstacle density. 24 trials each.",
         size=13, color=MUTED)

# Main chart (left)
chart_path = os.path.join(FIG, "results_combined.png")
if os.path.exists(chart_path):
    s.shapes.add_picture(chart_path, Inches(0.4), Inches(1.9),
                         width=Inches(7.9))

# Right: Key findings card
kx = 8.55
kw = 4.3
add_rect(s, kx, 1.9, kw, 5.0, fill=CARD)
add_text(s, kx + 0.25, 2.05, kw - 0.5, 0.45, "Key Findings",
         size=16, bold=True, color=ACCENT)

# Finding 1
add_dot(s, kx + 0.3, 2.72, d=0.22, color=AST)
add_text(s, kx + 0.65, 2.62, kw - 1.0, 0.35,
         "A* halves Dijkstra", size=13, bold=True, color=WHITE)
add_text(s, kx + 0.65, 2.95, kw - 1.0, 0.8,
         "Around 50 percent of Dijkstra's expansions at every density. Manhattan heuristic pays off consistently.",
         size=10, color=MUTED)

# Finding 2
add_dot(s, kx + 0.3, 3.87, d=0.22, color=JPS)
add_text(s, kx + 0.65, 3.77, kw - 1.0, 0.35,
         "JPS4 wins on open maps", size=13, bold=True, color=WHITE)
add_text(s, kx + 0.65, 4.10, kw - 1.0, 0.8,
         "At 15 percent density, JPS4 visits only 488 nodes versus 2062 for Dijkstra. A 76 percent reduction.",
         size=10, color=MUTED)

# Finding 3
add_dot(s, kx + 0.3, 5.02, d=0.22, color=WARN)
add_text(s, kx + 0.65, 4.92, kw - 1.0, 0.35,
         "JPS4 loses on dense maps", size=13, bold=True, color=WHITE)
add_text(s, kx + 0.65, 5.25, kw - 1.0, 0.8,
         "At 28 percent, JPS4 expands 970 cells versus A*'s 629. Shorter corridors mean shorter jumps.",
         size=10, color=MUTED)

# Correctness footer inside card
add_rect(s, kx + 0.25, 6.1, kw - 0.5, 0.7, fill=PANEL)
add_text(s, kx + 0.25, 6.1, kw - 0.5, 0.35,
         "CORRECTNESS", size=9, bold=True, color=JPS, align=PP_ALIGN.CENTER)
add_text(s, kx + 0.25, 6.4, kw - 0.5, 0.35,
         "All three return equal path lengths on every trial.",
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_footer(s, 4)

# =============================================================================
# Save
# =============================================================================
out_path = os.path.join(HERE, "ADA_Interim_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Size : {os.path.getsize(out_path)} bytes")
