"""
Generate ADA Interim Presentation — Methodology & Results Focus
Run: python make_presentation.py
Output: ADA_Interim_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colours ───────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0F, 0x17, 0x2A)
ACCENT     = RGBColor(0x3B, 0x82, 0xF6)   # blue   — JPS4
ACCENT2    = RGBColor(0x10, 0xB9, 0x81)   # teal   — A*
ACCENT3    = RGBColor(0xF5, 0x9E, 0x0B)   # amber  — Dijkstra
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY   = RGBColor(0x64, 0x74, 0x8B)
CARD_BG    = RGBColor(0x1E, 0x29, 0x3B)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def box(slide, l, t, w, h, bg=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def multiline(slide, lines, l, t, w, h):
    """lines: list of (text, size, bold, color, align, italic, space_before_pt)"""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, color, align, italic, sp) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if sp:
            p.space_before = Pt(sp)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color

def topbar(slide, color=ACCENT):
    b = slide.shapes.add_shape(1, Inches(0), Inches(0.55), W, Inches(0.045))
    b.fill.solid(); b.fill.fore_color.rgb = color
    b.line.fill.background()

def slide_num(slide, n, total=5):
    txt(slide, f"{n} / {total}", 12.5, 7.1, 0.8, 0.3,
        size=10, color=MID_GRAY, align=PP_ALIGN.RIGHT)

def heading(slide, text):
    txt(slide, text, 0.4, 0.1, 12.2, 0.55, size=26, bold=True, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_bg(s1)

# Left accent strip
strip = s1.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), H)
strip.fill.solid(); strip.fill.fore_color.rgb = ACCENT
strip.line.fill.background()

txt(s1, "ADA INTERIM DEMO  ·  APRIL 2026",
    0.8, 1.5, 11, 0.4, size=11, bold=True, color=ACCENT)

multiline(s1, [
    ("Benchmarking Shortest-Path", 36, True, WHITE, PP_ALIGN.LEFT, False, 0),
    ("Algorithms on 4-Connected Grid Maps", 36, True, WHITE, PP_ALIGN.LEFT, False, 2),
], 0.8, 2.0, 11.5, 1.8)

txt(s1, "Dijkstra  ·  A*  ·  Jump Point Search (JPS4)",
    0.8, 3.85, 10, 0.5, size=18, italic=True, color=ACCENT2)

bar = s1.shapes.add_shape(1, Inches(0.8), Inches(4.5), Inches(4.5), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = MID_GRAY
bar.line.fill.background()

multiline(s1, [
    ("Jazib Waqas   ·   Salman Adnan   ·   Hunain Abbas", 14, True, LIGHT_GRAY, PP_ALIGN.LEFT, False, 0),
    ("jw08048  ·  sa07885  ·  sh08466  ·  Habib University", 11, False, MID_GRAY, PP_ALIGN.LEFT, False, 5),
], 0.8, 4.65, 9, 1.0)

slide_num(s1, 1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — METHODOLOGY
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_bg(s2)
topbar(s2)
heading(s2, "Methodology — How We Built the Benchmark")

# ── Left: Benchmark design ────────────────────────────────────────────────────
txt(s2, "BENCHMARK DESIGN", 0.4, 0.72, 6.0, 0.32, size=10, bold=True, color=ACCENT)

multiline(s2, [
    ("Grid setup",  13, True,  ACCENT2,    PP_ALIGN.LEFT, False, 0),
    ("Random 50×50 grids with obstacles at 3 densities: 15%, 28%, 40%",
                   12, False, LIGHT_GRAY,  PP_ALIGN.LEFT, False, 4),
    ("Start and goal placed as far apart as possible; trials only counted when all 3 algorithms found a valid path",
                   12, False, LIGHT_GRAY,  PP_ALIGN.LEFT, False, 4),
    ("Fixed random seed (42), 24 trials per density for stable means",
                   12, False, LIGHT_GRAY,  PP_ALIGN.LEFT, False, 4),

    ("What we measure", 13, True,  ACCENT2, PP_ALIGN.LEFT, False, 12),
    ("Node expansions — how many cells each algorithm processed (hardware-independent)",
                   12, False, LIGHT_GRAY,  PP_ALIGN.LEFT, False, 4),
    ("Wall-clock time in ms — includes all per-step overhead",
                   12, False, LIGHT_GRAY,  PP_ALIGN.LEFT, False, 4),
    ("Path length — confirms all 3 return the same shortest path (correctness check)",
                   12, False, LIGHT_GRAY,  PP_ALIGN.LEFT, False, 4),
], 0.4, 1.1, 6.1, 5.6)

# ── Right: Shared loop design ─────────────────────────────────────────────────
txt(s2, "WHY THE COMPARISON IS FAIR", 6.7, 0.72, 6.2, 0.32, size=10, bold=True, color=ACCENT3)

box(s2, 6.65, 1.05, 6.3, 3.5, bg=CARD_BG)

multiline(s2, [
    ("One search loop, three plug-ins", 14, True, ACCENT3, PP_ALIGN.LEFT, False, 0),
    ("All three algorithms run on the same underlying loop:",
                   12, False, LIGHT_GRAY, PP_ALIGN.LEFT, False, 6),
], 6.85, 1.2, 5.9, 0.9)

plug_items = [
    (ACCENT3, "Dijkstra",  "heuristic = 0,  neighbours = 4 adjacent cells"),
    (ACCENT2, "A*",        "heuristic = Manhattan distance,  neighbours = 4 adjacent cells"),
    (ACCENT,  "JPS4",      "heuristic = Manhattan distance,  neighbours = jump points only"),
]
for i, (col, name, desc) in enumerate(plug_items):
    py = 2.2 + i * 0.85
    dot = s2.shapes.add_shape(9, Inches(6.9), Inches(py + 0.09), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = col
    dot.line.fill.background()
    multiline(s2, [
        (name, 12, True, col, PP_ALIGN.LEFT, False, 0),
        (desc, 11, False, LIGHT_GRAY, PP_ALIGN.LEFT, False, 2),
    ], 7.2, py, 5.5, 0.78)

txt(s2, "Any difference in results is purely algorithmic — not a coding artefact.",
    6.85, 4.62, 5.9, 0.4, size=11, italic=True, color=ACCENT2)

# ── Bottom: density rationale ─────────────────────────────────────────────────
box(s2, 6.65, 5.0, 6.3, 1.75, bg=CARD_BG)
multiline(s2, [
    ("Why three densities?", 13, True, ACCENT, PP_ALIGN.LEFT, False, 0),
    ("Low (15%): open map — long corridors, few forced choices", 11, False, LIGHT_GRAY, PP_ALIGN.LEFT, False, 5),
    ("Medium (28%): mixed — obstacles start breaking up corridors", 11, False, LIGHT_GRAY, PP_ALIGN.LEFT, False, 4),
    ("High (40%): dense — many dead-ends, short routes between obstacles", 11, False, LIGHT_GRAY, PP_ALIGN.LEFT, False, 4),
], 6.85, 5.15, 5.9, 1.5)

slide_num(s2, 2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — RESULTS (numbers + figure)
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_bg(s3)
topbar(s3)
heading(s3, "Results — Node Expansions & Timing")

# Table header
box(s3, 0.35, 0.72, 6.1, 0.4, bg=ACCENT)
multiline(s3, [
    ("                   15% blocked        28% blocked        40% blocked", 10, True, WHITE, PP_ALIGN.LEFT, False, 0),
], 0.38, 0.78, 6.0, 0.28)
box(s3, 0.35, 1.12, 6.1, 0.32, bg=RGBColor(0x1a,0x30,0x50))
multiline(s3, [
    ("Algorithm       Exp.    ms        Exp.    ms        Exp.    ms  ", 10, True, LIGHT_GRAY, PP_ALIGN.LEFT, False, 0),
], 0.38, 1.17, 6.0, 0.25)

rows = [
    ("Dijkstra", "2062", "7.07", "1659", "4.96",  "915", "2.73", ACCENT3, CARD_BG),
    ("A*",       "1020", "3.74",  "629", "2.11",  "513", "1.70", ACCENT2, RGBColor(0x18,0x26,0x34)),
    ("JPS4",      "488", "4.43",  "970", "6.22",  "682", "3.13", ACCENT,  CARD_BG),
]
for ri, (alg, e1, t1, e2, t2, e3, t3, col, bg) in enumerate(rows):
    ry = 1.44 + ri * 0.46
    box(s3, 0.35, ry, 6.1, 0.46, bg=bg)
    row_str = f"{alg:<10}  {e1:>6}  {t1:>5}     {e2:>6}  {t2:>5}     {e3:>6}  {t3:>5}"
    txt(s3, row_str, 0.5, ry + 0.1, 5.9, 0.3, size=11, bold=(alg == "JPS4"), color=col)

txt(s3, "Exp = mean node expansions   ·   ms = mean wall-clock time   ·   24 trials, 50×50 grid",
    0.35, 2.82, 6.1, 0.3, size=8.5, italic=True, color=MID_GRAY)

# Figure
fig_candidates = [
    os.path.join(os.path.dirname(__file__), "results_combined.png"),
    os.path.join(os.path.dirname(__file__), "..", "results_combined.png"),
    os.path.join(os.path.dirname(__file__), "..", "results", "results_combined.png"),
]
fig_path = None
for p in fig_candidates:
    if os.path.exists(p):
        fig_path = p
        break

if fig_path:
    s3.shapes.add_picture(fig_path, Inches(0.35), Inches(3.1), Inches(6.1), Inches(3.5))
else:
    box(s3, 0.35, 3.1, 6.1, 3.5, bg=CARD_BG)
    txt(s3, "[Insert results_combined.png here]",
        0.5, 4.5, 5.8, 0.5, size=11, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

txt(s3, "Figure 1: Mean expansions (left) and wall-clock time in ms (right) vs. obstacle density",
    0.35, 6.65, 6.1, 0.32, size=8.5, italic=True, color=MID_GRAY)

# Right: observations
txt(s3, "WHAT STANDS OUT", 6.7, 0.72, 6.2, 0.32, size=10, bold=True, color=ACCENT)

obs = [
    (ACCENT2, "A* vs Dijkstra",
     "A* consistently uses ~50% fewer node expansions than Dijkstra at every density tested. The Manhattan distance heuristic reliably cuts out the wrong half of the search space."),
    (ACCENT, "JPS4 at low density (15%)",
     "JPS4 expands only 488 nodes vs A*'s 1020 — a 52% reduction. Long open corridors let JPS4 jump far, visiting only the cells that actually matter."),
    (ACCENT3, "JPS4 at medium density (28%)",
     "JPS4's count rises to 970 — above A*'s 629. Obstacles shorten corridors, so each jump covers less ground. The pruning still runs, but returns less value."),
    (MID_GRAY, "Timing inversion",
     "JPS4 is slower on the clock despite fewer expansions. Its inner loop is more complex than A*'s 4-way lookup. In Python this cost dominates. In compiled code (C++/Java) it disappears — confirmed by the original literature."),
]

for oi, (col, title, body) in enumerate(obs):
    oy = 1.05 + oi * 1.52
    box(s3, 6.65, oy, 6.35, 1.42, bg=CARD_BG)
    edge = s3.shapes.add_shape(1, Inches(6.65), Inches(oy), Inches(0.06), Inches(1.42))
    edge.fill.solid(); edge.fill.fore_color.rgb = col
    edge.line.fill.background()
    multiline(s3, [
        (title, 12, True,  col,        PP_ALIGN.LEFT, False, 0),
        (body,  10, False, LIGHT_GRAY, PP_ALIGN.LEFT, False, 5),
    ], 6.85, oy + 0.12, 5.95, 1.2)

slide_num(s3, 3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DISCUSSION (what the results mean)
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_bg(s4)
topbar(s4, ACCENT2)
heading(s4, "Discussion — What the Results Tell Us")

# ── Main takeaway ─────────────────────────────────────────────────────────────
box(s4, 0.35, 0.7, 12.63, 0.88, bg=CARD_BG)
multiline(s4, [
    ("Core finding:", 13, True, ACCENT2, PP_ALIGN.LEFT, False, 0),
    ("JPS4's advantage is real, but conditional. It dominates on open maps; it loses its edge as obstacles fragment the space. The cross-over point is the key open question.", 13, False, LIGHT_GRAY, PP_ALIGN.LEFT, False, 5),
], 0.55, 0.82, 12.2, 0.68)

# ── Three insight cards ───────────────────────────────────────────────────────
cards = [
    (ACCENT, "Why JPS4 wins on open maps",
     [
         "In a 4-connected grid, moving in a straight line through empty cells is always symmetric — every intermediate cell has identical cost and the same 3 neighbours.",
         "JPS4 recognises this and skips directly to the next forced choice: a cell where an obstacle means you cannot continue straight without losing options.",
         "Fewer cells in the open list = faster convergence. At 15% density this produces a 76% reduction in expansions compared to Dijkstra.",
     ]),
    (ACCENT3, "Why JPS4 loses at medium density",
     [
         "Obstacles break corridors into short segments. Each jump now covers only 3–5 cells before hitting a wall, so the savings from skipping decrease.",
         "But the cost of checking for jump points on each step stays constant. The algorithm does more work per expansion than A* and gets less payoff.",
         "This is the density cross-over: the point where the overhead starts outweighing the pruning benefit. We have not yet pinned it down exactly.",
     ]),
    (MID_GRAY, "What the timing result means",
     [
         "JPS4 is slower in ms even where it expands fewer nodes. This is a Python prototype issue — the interpreter cost of the more complex inner loop dominates.",
         "In C++ and Java implementations the per-step overhead is negligible (as shown in the original Harabor & Grastien 2011 results).",
         "Node expansions are the fair algorithmic metric at this stage. Timing becomes meaningful once we move to compiled benchmarks.",
     ]),
]

for ci, (col, title, bullets) in enumerate(cards):
    cx = 0.35 + ci * 4.35
    box(s4, cx, 1.72, 4.15, 4.8, bg=CARD_BG)
    topb = s4.shapes.add_shape(1, Inches(cx), Inches(1.72), Inches(4.15), Inches(0.25))
    topb.fill.solid(); topb.fill.fore_color.rgb = col
    topb.line.fill.background()
    txt(s4, title, cx + 0.12, 1.76, 3.9, 0.28, size=11, bold=True, color=WHITE)
    lines = [(("• " + b), 10, False, LIGHT_GRAY, PP_ALIGN.LEFT, False, 6 if j>0 else 0)
             for j, b in enumerate(bullets)]
    multiline(s4, lines, cx + 0.12, 2.05, 3.9, 4.35)

# ── Bottom: open question ─────────────────────────────────────────────────────
box(s4, 0.35, 6.65, 12.63, 0.62, bg=RGBColor(0x0A, 0x0F, 0x1E))
multiline(s4, [
    ("Open question for Week 13:", 12, True, ACCENT, PP_ALIGN.LEFT, False, 0),
    ("At what obstacle density does A* become more efficient than JPS4 on a 4-connected grid? Sweeping 5–50% in 5% steps will give us the full performance curve.", 12, False, LIGHT_GRAY, PP_ALIGN.LEFT, False, 0),
], 0.6, 6.75, 12.1, 0.48)

slide_num(s4, 4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DEMO PREVIEW & NEXT STEPS (brief, not the main event)
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_bg(s5)
topbar(s5, ACCENT3)
heading(s5, "What's Next — Demo & Final Benchmarks")

# Left: snake game description
txt(s5, "SNAKE-GAME VISUALIZATION (IN PROGRESS)", 0.4, 0.72, 6.2, 0.32, size=10, bold=True, color=ACCENT3)
box(s5, 0.35, 1.05, 6.15, 4.55, bg=CARD_BG)

multiline(s5, [
    ("Why a Snake game?", 14, True, ACCENT3, PP_ALIGN.LEFT, False, 0),
    ("The report gives the numbers. The demo makes them visible.", 12, False, LIGHT_GRAY, PP_ALIGN.LEFT, False, 5),
    ("A snake navigates toward food on a grid. Each algorithm runs separately, and the cells it explores are highlighted in real time:", 12, False, LIGHT_GRAY, PP_ALIGN.LEFT, False, 8),
], 0.55, 1.2, 5.75, 1.5)

demo_rows = [
    (ACCENT3, "Dijkstra",  "Wide fan-out — explores in every direction before finding the goal"),
    (ACCENT2, "A*",        "Directed cone — focuses straight toward the food"),
    (ACCENT,  "JPS4",      "Sparse jumps — long straight lines with only a few stops"),
]
for di, (col, name, desc) in enumerate(demo_rows):
    dy = 2.85 + di * 0.88
    dot = s5.shapes.add_shape(9, Inches(0.55), Inches(dy+0.12), Inches(0.2), Inches(0.2))
    dot.fill.solid(); dot.fill.fore_color.rgb = col
    dot.line.fill.background()
    multiline(s5, [
        (name, 12, True, col, PP_ALIGN.LEFT, False, 0),
        (desc, 11, False, LIGHT_GRAY, PP_ALIGN.LEFT, False, 3),
    ], 0.88, dy, 5.5, 0.82)

txt(s5, "The visual difference makes the algorithmic trade-offs immediately intuitive for any audience.",
    0.55, 5.55, 5.8, 0.4, size=11, italic=True, color=ACCENT2)

# ── Right: timeline ───────────────────────────────────────────────────────────
txt(s5, "REMAINING TIMELINE", 6.7, 0.72, 6.2, 0.32, size=10, bold=True, color=ACCENT)

timeline = [
    (ACCENT,  "Week 13 — Scale Up",
     "100×100 and 200×200 grids. Density sweep 5–50% to find the A*/JPS4 cross-over. Complete the Snake-game demo."),
    (ACCENT2, "Week 14 — Final Analysis",
     "Finalize all plots and tables. Write the discussion comparing results to literature benchmarks. Polish the paper."),
    (ACCENT3, "Week 15 — Presentation",
     "Prepare final slides, rehearse the live demo, and submit by deadline."),
]

for ti, (col, week, desc) in enumerate(timeline):
    ty = 1.05 + ti * 2.1
    box(s5, 6.65, ty, 6.3, 1.92, bg=CARD_BG)
    edge = s5.shapes.add_shape(1, Inches(6.65), Inches(ty), Inches(0.07), Inches(1.92))
    edge.fill.solid(); edge.fill.fore_color.rgb = col
    edge.line.fill.background()
    multiline(s5, [
        (week, 13, True, col, PP_ALIGN.LEFT, False, 0),
        (desc, 11, False, LIGHT_GRAY, PP_ALIGN.LEFT, False, 6),
    ], 6.88, ty + 0.18, 5.9, 1.62)

slide_num(s5, 5)

# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "ADA_Interim_Presentation_v2.pptx")
prs.save(out)
print(f"✅ Saved: {out}")
